"""Generate PikoGPT PowerPoint presentation (English version)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BG_DARK = RGBColor(0x12, 0x12, 0x20)
BG_CARD = RGBColor(0x1C, 0x1C, 0x30)
PURPLE_MAIN = RGBColor(0x81, 0x8C, 0xF8)
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_DARK = RGBColor(0x63, 0x66, 0xF1)
GREEN = RGBColor(0x34, 0xD3, 0x99)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
WHITE = RGBColor(0xE8, 0xE8, 0xF0)
GRAY = RGBColor(0xA5, 0xA5, 0xC0)
DARK_GRAY = RGBColor(0x6B, 0x6B, 0x8A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, highlight_color=PURPLE_LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)

        if "|" in item:
            parts = item.split("|", 1)
            run1 = p.add_run()
            run1.text = "▸ " + parts[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = highlight_color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = "▸ " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return txBox


def add_slide_number(slide, num, total=12):
    add_text_box(slide, Inches(12.3), Inches(0.3), Inches(0.8), Inches(0.4),
                 f"{num:02d}", font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT,
                 font_name="Consolas")


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.5), color=PURPLE_MAIN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = PURPLE_MAIN
bar.line.fill.background()

add_text_box(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(1.5),
             "PikoGPT", font_size=72, color=PURPLE_MAIN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(3.3), SLIDE_W, Inches(0.7),
             "A Small Language Model — Built from Scratch", font_size=24, color=GRAY,
             alignment=PP_ALIGN.CENTER)

badge = add_shape_bg(slide, Inches(5.2), Inches(4.4), Inches(2.9), Inches(0.55), BG_CARD, 0.5)
badge.line.color.rgb = RGBColor(0x40, 0x44, 0x80)
badge.line.width = Pt(1)
add_text_box(slide, Inches(5.2), Inches(4.42), Inches(2.9), Inches(0.55),
             "TEAM FUNKYAI", font_size=16, color=PURPLE_MAIN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.4), SLIDE_W, Inches(0.5),
             "NLP with LLMs  ·  University of St. Gallen  ·  Spring 2026", font_size=14,
             color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.1), SLIDE_W, Inches(0.4),
             "Filipp  ·  Arabella  ·  Roman", font_size=15,
             color=GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 2: AGENDA
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 2)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.8),
             "Agenda", font_size=40, color=PURPLE_MAIN, bold=True)

left_items = [
    "01  Project Goal & Motivation",
    "02  Team & Responsibilities",
    "03  Data Pipeline",
    "04  Model Architecture",
    "05  Model Configurations",
]
right_items = [
    "06  Training & Optimization",
    "07  Hyperparameter Tuning",
    "08  Inference & Text Generation",
    "09  Technology Stack",
    "10  Summary & Outlook",
]

for i, item in enumerate(left_items):
    y = Inches(1.8 + i * 0.7)
    add_accent_bar(slide, Inches(1.0), y, height=Inches(0.35))
    num, rest = item.split("  ", 1)
    txBox = slide.shapes.add_textbox(Inches(1.2), y, Inches(5), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = num + "   "
    r1.font.size = Pt(20)
    r1.font.color.rgb = PURPLE_LIGHT
    r1.font.bold = True
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(20)
    r2.font.color.rgb = WHITE
    r2.font.name = "Calibri"

for i, item in enumerate(right_items):
    y = Inches(1.8 + i * 0.7)
    add_accent_bar(slide, Inches(7.0), y, height=Inches(0.35))
    num, rest = item.split("  ", 1)
    txBox = slide.shapes.add_textbox(Inches(7.2), y, Inches(5), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = num + "   "
    r1.font.size = Pt(20)
    r1.font.color.rgb = PURPLE_LIGHT
    r1.font.bold = True
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(20)
    r2.font.color.rgb = WHITE
    r2.font.name = "Calibri"


# =====================================================================
# SLIDE 3: PROJECT GOAL & MOTIVATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 3)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Project Goal & Motivation", font_size=40, color=PURPLE_MAIN, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(3), Inches(0.5),
             "What is PikoGPT?", font_size=22, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(2.5), [
    "Decoder-only Transformer| language model",
    "Trained on |OpenWebText data",
    "Model size: |10–40M parameters",
    "Fully |implemented from scratch",
], font_size=17)

add_text_box(slide, Inches(0.8), Inches(4.4), Inches(3), Inches(0.5),
             "Project Goal", font_size=22, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(2), [
    "In 12 weeks| from zero to a working LLM",
    "Full pipeline:| Data → Training → Chat Interface",
    "Deepen understanding| of LLM fundamentals",
], font_size=17)

# Stats boxes
stats = [("37M", "Parameters\n(Large)"), ("500K", "Training\nDocuments"),
         ("512", "Context\nTokens"), ("12", "Week\nProject")]
for i, (num, label) in enumerate(stats):
    col = i % 2
    row = i // 2
    x = Inches(8.0 + col * 2.4)
    y = Inches(1.6 + row * 2.6)
    card = add_shape_bg(slide, x, y, Inches(2.1), Inches(2.1), BG_CARD, 0.08)
    card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
    card.line.width = Pt(0.75)
    add_text_box(slide, x, y + Inches(0.3), Inches(2.1), Inches(0.8),
                 num, font_size=44, color=PURPLE_MAIN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.3), Inches(2.1), Inches(0.7),
                 label, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 4: TEAM
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 4)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Team FunkyAI", font_size=40, color=PURPLE_MAIN, bold=True)

team = [
    ("Filipp", "Data Engineer", PURPLE_MAIN, [
        "Data pipeline & quality assurance",
        "Benchmarks (LAMBADA, HellaSwag, WinoGrande)",
        "Leaderboard evaluations & submissions",
        "Results analysis & visualizations",
    ]),
    ("Arabella", "ML Engineer", PURPLE_LIGHT, [
        "Model architecture & configuration",
        "Hyperparameter tuning with Optuna",
        "GPU training runs & monitoring",
        "Ablation studies & model selection",
    ]),
    ("Roman", "Platform Engineer", GREEN, [
        "Chat interface (Gradio / Streamlit)",
        "Documentation & README",
        "Code quality & testing",
        "Infrastructure & DevOps",
    ]),
]

for i, (name, role, role_color, tasks) in enumerate(team):
    x = Inches(0.8 + i * 4.0)
    card = add_shape_bg(slide, x, Inches(1.8), Inches(3.6), Inches(4.8), BG_CARD, 0.04)
    card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
    card.line.width = Pt(0.75)
    add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3), Inches(0.5),
                 name, font_size=24, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(3), Inches(0.4),
                 role, font_size=16, color=role_color, bold=True)
    add_bullet_list(slide, x + Inches(0.3), Inches(3.3), Inches(3.0), Inches(3.0),
                    tasks, font_size=14, color=GRAY)


# =====================================================================
# SLIDE 5: DATA PIPELINE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 5)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Data Pipeline", font_size=40, color=PURPLE_MAIN, bold=True)

steps = ["📥 Stream\nOpenWebText", "🌐 Language\nDetection", "🔒 Test Set\nFiltering",
         "🧹 Remove HTML\nURLs, Code", "✅ Quality\nFilter", "💾 Save\nDataset"]

for i, step in enumerate(steps):
    x = Inches(0.5 + i * 2.1)
    card = add_shape_bg(slide, x, Inches(1.6), Inches(1.8), Inches(1.2), BG_CARD, 0.08)
    card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
    card.line.width = Pt(0.75)
    add_text_box(slide, x, Inches(1.7), Inches(1.8), Inches(1.1),
                 step, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.8), Inches(1.9), Inches(0.3), Inches(0.5),
                     "→", font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Left card
card_left = add_shape_bg(slide, Inches(0.8), Inches(3.3), Inches(5.6), Inches(3.7), BG_CARD, 0.03)
card_left.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card_left.line.width = Pt(0.75)

add_text_box(slide, Inches(1.1), Inches(3.5), Inches(5), Inches(0.4),
             "Preprocessing Steps", font_size=20, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(4.0), Inches(5.2), Inches(2.8), [
    "Language filter: |English-only texts (langdetect)",
    "Data leakage protection: |MD5 hashing of all test sentences",
    "HTML cleaning: |Remove tags and special characters",
    "URL filtering: |Remove HTTP(S) and www links",
    "Code blocks: |Remove markdown code blocks",
    "Quality check: |Min. 100 chars, no corruption",
], font_size=15)

# Right card
card_right = add_shape_bg(slide, Inches(6.8), Inches(3.3), Inches(5.7), Inches(3.7), BG_CARD, 0.03)
card_right.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card_right.line.width = Pt(0.75)

add_text_box(slide, Inches(7.1), Inches(3.5), Inches(5), Inches(0.4),
             "Usage", font_size=20, color=PURPLE_LIGHT, bold=True)

code_bg = add_shape_bg(slide, Inches(7.1), Inches(4.1), Inches(5.1), Inches(1.6),
                        RGBColor(0x0A, 0x0A, 0x18), 0.03)
code_bg.line.color.rgb = RGBColor(0x25, 0x28, 0x45)
code_bg.line.width = Pt(0.75)

add_text_box(slide, Inches(7.3), Inches(4.2), Inches(4.8), Inches(1.4),
             '# Run preprocessing\npython main.py \\\n  --stage preprocess \\\n  --num-samples 100000 \\\n  --output-path "data/processed/openwebtext_clean"',
             font_size=13, color=GREEN, font_name="Consolas")


# =====================================================================
# SLIDE 6: MODEL ARCHITECTURE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 6)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Model Architecture", font_size=40, color=PURPLE_MAIN, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
             "CausalTransformerLM", font_size=22, color=PURPLE_LIGHT, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.8),
             "Decoder-only transformer with causal masking,\nbased on PyTorch's nn.TransformerEncoder.",
             font_size=16, color=GRAY)

add_bullet_list(slide, Inches(0.8), Inches(3.1), Inches(6), Inches(3.5), [
    "Token Embedding| + learnable position embedding",
    "Multi-Head Self-Attention| with causal mask",
    "Feed-Forward Network| with GELU activation (4× d_model)",
    "Layer Normalization| + dropout per block",
    "Weight Tying:| LM head shares embedding weights",
    "Weight Initialization:| N(0, 0.02)",
    "Loss:| Cross-entropy on shifted labels",
], font_size=16)

# Architecture stack
arch_blocks = [
    ("LM Head (Linear → Vocab)", AMBER, RGBColor(0x3A, 0x2E, 0x10)),
    ("Final LayerNorm", GREEN, RGBColor(0x12, 0x30, 0x25)),
    ("Transformer Block × N\nSelf-Attention + FFN + LayerNorm", PURPLE_LIGHT, RGBColor(0x28, 0x1A, 0x3A)),
    ("Dropout", PURPLE_MAIN, RGBColor(0x1A, 0x1C, 0x38)),
    ("Token Emb + Position Emb", PURPLE_MAIN, RGBColor(0x1A, 0x1C, 0x38)),
]

for i, (text, text_color, bg_color) in enumerate(arch_blocks):
    y = Inches(1.5 + i * 1.15)
    block = add_shape_bg(slide, Inches(8.0), y, Inches(4.3), Inches(0.85), bg_color, 0.06)
    block.line.color.rgb = RGBColor(0x40, 0x44, 0x70)
    block.line.width = Pt(0.75)
    add_text_box(slide, Inches(8.0), y + Inches(0.05), Inches(4.3), Inches(0.75),
                 text, font_size=14, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(arch_blocks) - 1:
        add_text_box(slide, Inches(9.8), y + Inches(0.85), Inches(0.7), Inches(0.3),
                     "↑", font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(8.0), Inches(1.5 + len(arch_blocks) * 1.15), Inches(4.3), Inches(0.4),
             "Input Token IDs", font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 7: MODEL CONFIGURATIONS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 7)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Model Configurations", font_size=40, color=PURPLE_MAIN, bold=True)

headers = ["Configuration", "Parameters", "n_embd", "n_layer", "n_head", "Context", "Focus"]
col_widths = [2.5, 1.3, 1.2, 1.2, 1.2, 1.2, 3.2]

header_y = Inches(1.8)
header_bg = add_shape_bg(slide, Inches(0.6), header_y, Inches(11.8), Inches(0.5),
                          RGBColor(0x1A, 0x1C, 0x30))
header_bg.line.fill.background()

x_pos = 0.6
for j, (header, w) in enumerate(zip(headers, col_widths)):
    add_text_box(slide, Inches(x_pos), header_y + Inches(0.05), Inches(w), Inches(0.4),
                 header, font_size=12, color=PURPLE_MAIN, bold=True, font_name="Calibri")
    x_pos += w

rows = [
    ["Default  [Test]", "~16M", "256", "4", "4", "128", "Quick testing"],
    ["Large  [⭐ Primary]", "~37M", "384", "10", "6", "512", "Depth/width balance"],
    ["Deep  [Variant]", "~34M", "320", "14", "8", "512", "Deep pattern learning"],
    ["Full Context  [Variant]", "~33M", "352", "10", "8", "1024", "Long-range dependencies"],
]

for i, row in enumerate(rows):
    y = Inches(2.5 + i * 0.75)
    if i % 2 == 0:
        row_bg = add_shape_bg(slide, Inches(0.6), y, Inches(11.8), Inches(0.65),
                               RGBColor(0x16, 0x16, 0x28))
        row_bg.line.fill.background()
    x_pos = 0.6
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        color = WHITE if j == 0 else GRAY
        add_text_box(slide, Inches(x_pos), y + Inches(0.1), Inches(w), Inches(0.5),
                     cell, font_size=15, color=color, font_name="Calibri")
        x_pos += w

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
             "All models use the GPT-2 Tokenizer (vocab: 50,257), AdamW optimizer (weight decay 0.01), and gradient clipping (norm 1.0).",
             font_size=14, color=DARK_GRAY)


# =====================================================================
# SLIDE 8: TRAINING & OPTIMIZATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 8)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Training & Optimization", font_size=40, color=PURPLE_MAIN, bold=True)

card = add_shape_bg(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2), BG_CARD, 0.03)
card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card.line.width = Pt(0.75)

add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.4),
             "Training Configuration", font_size=20, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(2.5), Inches(5.0), Inches(4), [
    "Optimizer: |AdamW (lr=3e-4, weight_decay=0.01)",
    "LR Schedule: |Cosine decay with linear warmup",
    "Warmup: |500 steps linear increase",
    "Gradient Clipping: |Max norm 1.0",
    "Gradient Accumulation: |2–4 steps",
    "Batch Size: |8–16 (effective up to 64)",
    "Max Steps: |50,000",
    "torch.compile: |Enabled for performance",
], font_size=15)

card2 = add_shape_bg(slide, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.5), BG_CARD, 0.03)
card2.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card2.line.width = Pt(0.75)

add_text_box(slide, Inches(7.1), Inches(1.9), Inches(5), Inches(0.4),
             "Distributed Training", font_size=20, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(7.1), Inches(2.5), Inches(5.0), Inches(1.6), [
    "DDP: |PyTorch DistributedDataParallel",
    "Multi-GPU: |Scales across multiple GPUs",
    "Checkpointing: |Save every 500 steps",
    "Resume: |Continue training from checkpoint",
], font_size=15)

card3 = add_shape_bg(slide, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), BG_CARD, 0.03)
card3.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card3.line.width = Pt(0.75)

add_text_box(slide, Inches(7.1), Inches(4.7), Inches(5), Inches(0.4),
             "Monitoring", font_size=20, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(7.1), Inches(5.2), Inches(5.0), Inches(1.5), [
    "JSONL metrics| (train + eval per step)",
    "Gradient norm| monitoring & health checks",
    "Perplexity| tracking across epochs",
], font_size=15)


# =====================================================================
# SLIDE 9: HYPERPARAMETER TUNING
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 9)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Hyperparameter Tuning", font_size=40, color=PURPLE_MAIN, bold=True)

card = add_shape_bg(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2), BG_CARD, 0.03)
card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card.line.width = Pt(0.75)

add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.4),
             "Optuna-based Search", font_size=20, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(2.4), Inches(5.0), Inches(1.5), [
    "Framework: |Optuna with SQLite storage",
    "Dashboard: |Optuna Dashboard for visualization",
    "Budget: |Max 40M parameter limit",
    "Fast evaluation: |500 steps per trial",
], font_size=15)

add_text_box(slide, Inches(1.1), Inches(4.1), Inches(5), Inches(0.4),
             "Tuned Hyperparameters", font_size=18, color=PURPLE_LIGHT, bold=True)

hp_rows = [
    ("Learning Rate", "1e-5 – 1e-3"),
    ("Batch Size", "8 – 32"),
    ("n_layer", "4 – 14"),
    ("n_head", "4 – 8"),
    ("n_embd", "256 – 384"),
    ("Dropout", "0.0 – 0.3"),
    ("Warmup Steps", "100 – 1,000"),
    ("Weight Decay", "0.005 – 0.1"),
]

for i, (param, range_) in enumerate(hp_rows):
    y = Inches(4.6 + i * 0.3)
    add_text_box(slide, Inches(1.3), y, Inches(2.5), Inches(0.3),
                 param, font_size=13, color=WHITE)
    add_text_box(slide, Inches(3.8), y, Inches(2), Inches(0.3),
                 range_, font_size=13, color=PURPLE_LIGHT)

card2 = add_shape_bg(slide, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2), BG_CARD, 0.03)
card2.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card2.line.width = Pt(0.75)

add_text_box(slide, Inches(7.1), Inches(1.9), Inches(5), Inches(0.4),
             "Usage", font_size=20, color=PURPLE_LIGHT, bold=True)

code_bg2 = add_shape_bg(slide, Inches(7.1), Inches(2.5), Inches(5.1), Inches(2.0),
                          RGBColor(0x0A, 0x0A, 0x18), 0.03)
code_bg2.line.color.rgb = RGBColor(0x25, 0x28, 0x45)
code_bg2.line.width = Pt(0.75)

add_text_box(slide, Inches(7.3), Inches(2.6), Inches(4.8), Inches(1.8),
             '# Start hyperparameter search\npython -m src.tuning.optuna_search \\\n  --n-trials 20\n\n# Open dashboard\noptuna-dashboard \\\n  sqlite:///optuna_pikogpt.db \\\n  --host 127.0.0.1 --port 8080',
             font_size=12, color=GREEN, font_name="Consolas")

add_text_box(slide, Inches(7.1), Inches(4.8), Inches(5), Inches(0.4),
             "Fixed Parameters", font_size=18, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(7.1), Inches(5.3), Inches(5.0), Inches(1.2), [
    "Vocab Size: |50,257 (GPT-2 Tokenizer)",
    "Context Length: |256–512 tokens",
    "Activation: |GELU",
    "Optimizer: |AdamW (always)",
], font_size=15)


# =====================================================================
# SLIDE 10: INFERENCE & TEXT GENERATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 10)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Inference & Text Generation", font_size=40, color=PURPLE_MAIN, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5),
             "Generation Pipeline", font_size=22, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), [
    "Autoregressive decoding| — token by token",
    "Greedy decoding| (temperature = 0)",
    "Sampling| with adjustable temperature",
    "Sliding window| for texts beyond context length",
    "Leaderboard mode| for automated evaluation",
], font_size=16)

add_text_box(slide, Inches(0.8), Inches(4.6), Inches(3), Inches(0.4),
             "Usage", font_size=20, color=PURPLE_LIGHT, bold=True)

code_bg3 = add_shape_bg(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(1.8),
                          RGBColor(0x0A, 0x0A, 0x18), 0.03)
code_bg3.line.color.rgb = RGBColor(0x25, 0x28, 0x45)
code_bg3.line.width = Pt(0.75)

add_text_box(slide, Inches(1.0), Inches(5.2), Inches(5.2), Inches(1.6),
             'python main.py \\\n  --stage inference \\\n  --checkpoint model_final.pt \\\n  --prompt "The future of AI is" \\\n  --max-tokens 100 \\\n  --temperature 0.8',
             font_size=13, color=GREEN, font_name="Consolas")

card = add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.3), BG_CARD, 0.03)
card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card.line.width = Pt(0.75)

add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.4),
             "Checkpoint Format", font_size=20, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(7.3), Inches(2.3), Inches(5.0), Inches(1.5), [
    "state_dict| — model weights",
    "model| — architecture metadata (n_embd, n_layer, ...)",
    "tokenizer| — tokenizer name & configuration",
    "training_state| — optimizer, scheduler, step",
], font_size=15)

card2 = add_shape_bg(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(2.7), BG_CARD, 0.03)
card2.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card2.line.width = Pt(0.75)

add_text_box(slide, Inches(7.3), Inches(4.4), Inches(5), Inches(0.4),
             "Device Support", font_size=20, color=PURPLE_LIGHT, bold=True)

devices = [
    ("CUDA", "NVIDIA GPUs", PURPLE_MAIN),
    ("MPS", "Apple Silicon (M1/M2/M3/M4)", PURPLE_LIGHT),
    ("CPU", "Fallback for all systems", AMBER),
    ("Auto", "Automatic device detection", GREEN),
]

for i, (dev, desc, color) in enumerate(devices):
    y = Inches(5.0 + i * 0.4)
    add_text_box(slide, Inches(7.5), y, Inches(1.2), Inches(0.35),
                 dev, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(8.7), y, Inches(3.5), Inches(0.35),
                 desc, font_size=14, color=GRAY)


# =====================================================================
# SLIDE 11: TECHNOLOGY STACK
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 11)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             "Technology Stack", font_size=40, color=PURPLE_MAIN, bold=True)

tech_cats = [
    ("Deep Learning", [
        "PyTorch| — model & training",
        "torch.compile| — compilation optimization",
        "DDP| — distributed training",
        "Transformers| — GPT-2 tokenizer",
    ]),
    ("Data & Analysis", [
        "HuggingFace Datasets| — data handling",
        "langdetect| — language detection",
        "Pandas| — data analysis",
        "Matplotlib| — visualization",
    ]),
    ("Engineering", [
        "Pydantic v2| — config validation",
        "Optuna| — hyperparameter tuning",
        "TOML / YAML| — config files",
        "uv| — fast package manager",
    ]),
]

for i, (title, items) in enumerate(tech_cats):
    x = Inches(0.6 + i * 4.1)
    card = add_shape_bg(slide, x, Inches(1.6), Inches(3.8), Inches(3.0), BG_CARD, 0.03)
    card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
    card.line.width = Pt(0.75)
    add_text_box(slide, x + Inches(0.3), Inches(1.8), Inches(3.4), Inches(0.4),
                 title, font_size=18, color=PURPLE_LIGHT, bold=True)
    add_bullet_list(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(2.0),
                    items, font_size=14, color=GRAY)

add_text_box(slide, Inches(0.8), Inches(5.0), Inches(5), Inches(0.4),
             "Project Structure", font_size=20, color=PURPLE_LIGHT, bold=True)

struct_bg = add_shape_bg(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.6),
                          RGBColor(0x0A, 0x0A, 0x18), 0.02)
struct_bg.line.color.rgb = RGBColor(0x25, 0x28, 0x45)
struct_bg.line.width = Pt(0.75)

struct_text = (
    "nlp-pikogpt-funkyai/\n"
    "├── src/data/preprocessing.py      # Data pipeline\n"
    "├── src/training/{config,stage,utils}.py  # Training + model\n"
    "├── src/inference/stage.py         # Text generation\n"
    "├── src/tuning/optuna_search.py    # HP search\n"
    "├── configs/train_*.toml           # 4 configurations\n"
    "└── main.py                        # CLI entry point"
)
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11), Inches(1.4),
             struct_text, font_size=12, color=GREEN, font_name="Consolas")


# =====================================================================
# SLIDE 12: SUMMARY & OUTLOOK
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 12)

add_text_box(slide, Inches(0), Inches(0.5), SLIDE_W, Inches(0.8),
             "Summary & Outlook", font_size=40, color=PURPLE_MAIN, bold=True,
             alignment=PP_ALIGN.CENTER)

card = add_shape_bg(slide, Inches(1.0), Inches(1.6), Inches(5.3), Inches(4.8), BG_CARD, 0.03)
card.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card.line.width = Pt(0.75)

add_text_box(slide, Inches(1.3), Inches(1.8), Inches(4.8), Inches(0.5),
             "✅  Achieved", font_size=22, color=GREEN, bold=True)

add_bullet_list(slide, Inches(1.3), Inches(2.4), Inches(4.7), Inches(3.8), [
    "Complete data pipeline with leakage protection",
    "Decoder-only transformer built from scratch",
    "4 model configurations (16M–37M parameters)",
    "Distributed training with DDP",
    "Cosine LR scheduling with warmup",
    "Checkpoint system & resume support",
    "Inference pipeline with sampling",
    "Optuna hyperparameter tuning",
], font_size=15, color=GRAY)

card2 = add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.3), Inches(4.8), BG_CARD, 0.03)
card2.line.color.rgb = RGBColor(0x30, 0x33, 0x55)
card2.line.width = Pt(0.75)

add_text_box(slide, Inches(7.3), Inches(1.8), Inches(4.8), Inches(0.5),
             "🔮  Next Steps", font_size=22, color=PURPLE_LIGHT, bold=True)

add_bullet_list(slide, Inches(7.3), Inches(2.4), Inches(4.7), Inches(3.8), [
    "Full-scale GPU training",
    "Benchmark evaluation (LAMBADA, HellaSwag)",
    "Leaderboard submissions",
    "Chat interface (Gradio / Streamlit)",
    "Fine-tuning experiments",
    "Ablation studies",
    "Final report & poster",
], font_size=15, color=GRAY)

badge = add_shape_bg(slide, Inches(5.0), Inches(6.7), Inches(3.3), Inches(0.55), BG_CARD, 0.5)
badge.line.color.rgb = RGBColor(0x40, 0x44, 0x80)
badge.line.width = Pt(1)
add_text_box(slide, Inches(5.0), Inches(6.72), Inches(3.3), Inches(0.55),
             "Thank You! 🎉", font_size=18, color=PURPLE_MAIN, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SAVE
# =====================================================================
output_path = os.path.join(os.path.dirname(__file__), "PikoGPT_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
